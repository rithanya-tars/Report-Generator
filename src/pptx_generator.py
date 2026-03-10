"""
pptx_generator.py
Takes Claude's slide plan JSON and generates a polished PPTX using python-pptx.
Design: Clean, professional, Tars-branded with client logo support.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy


# ─── COLOR PALETTE ────────────────────────────────────────────────────────────
# Tars brand-inspired: deep purple primary, clean whites, accent highlights
C_TARS_PURPLE   = RGBColor(0x4A, 0x1E, 0x8A)   # Deep Tars purple
C_LIGHT_PURPLE  = RGBColor(0xEF, 0xE8, 0xFB)   # Very light purple bg
C_ACCENT        = RGBColor(0x7C, 0x3A, 0xED)   # Vibrant purple accent
C_WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT     = RGBColor(0x1A, 0x1A, 0x2E)   # Near black
C_MID_TEXT      = RGBColor(0x4A, 0x4A, 0x6A)   # Medium grey-blue
C_MUTED         = RGBColor(0x9A, 0x9A, 0xBA)   # Muted for captions
C_TABLE_HEADER  = RGBColor(0x4A, 0x1E, 0x8A)   # Purple table headers
C_TABLE_ALT     = RGBColor(0xF8, 0xF5, 0xFF)   # Light purple alt rows
C_GREEN         = RGBColor(0x10, 0xB9, 0x81)   # Positive metric
C_RED           = RGBColor(0xEF, 0x44, 0x44)   # Negative metric
C_PINK_BG       = RGBColor(0xFF, 0xF0, 0xF5)   # Slide title box bg (like Amex deck)

# Slide dimensions (16:9)
W = Inches(13.33)
H = Inches(7.5)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color, transparency=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    if transparency is not None:
        fill.fore_color.theme_color = None
    shape.line.fill.background()
    return shape


def set_text(tf, text, size, bold=False, color=C_DARK_TEXT, align=PP_ALIGN.LEFT, italic=False):
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Calibri"


def add_textbox(slide, text, x, y, w, h, size, bold=False, color=C_DARK_TEXT,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    set_text(tf, text, size, bold, color, align, italic)
    return txBox


def slide_background(slide, color=C_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_header(prs, slide, title_text, show_logos=True, assets=None):
    """Add the standard Tars + Client logo header and pink title box."""
    # Top bar
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.8), C_TARS_PURPLE)

    # Tars logo text (fallback if no logo image)
    if assets and assets.get("tars_logo") and os.path.exists(assets["tars_logo"]):
        try:
            slide.shapes.add_picture(assets["tars_logo"], Inches(0.2), Inches(0.1), height=Inches(0.6))
        except Exception:
            add_textbox(slide, "≡TARS", Inches(0.2), Inches(0.1), Inches(2), Inches(0.6),
                       22, bold=True, color=C_WHITE)
    else:
        add_textbox(slide, "≡TARS", Inches(0.2), Inches(0.1), Inches(2), Inches(0.6),
                   22, bold=True, color=C_WHITE)

    # Client logo
    if assets and assets.get("client_logo") and os.path.exists(assets["client_logo"]):
        try:
            slide.shapes.add_picture(assets["client_logo"], Inches(2.5), Inches(0.05),
                                     height=Inches(0.7))
        except Exception:
            pass

    # Pink title box (matching the Amex deck style)
    add_rect(slide, Inches(3.8), Inches(0.85), Inches(9.3), Inches(0.65), C_PINK_BG)
    # Left accent bar
    add_rect(slide, Inches(3.8), Inches(0.85), Inches(0.08), Inches(0.65), C_ACCENT)

    add_textbox(slide, title_text, Inches(4.0), Inches(0.88), Inches(9.0), Inches(0.6),
               20, bold=True, color=C_DARK_TEXT)


# ─── SLIDE BUILDERS ───────────────────────────────────────────────────────────

def build_cover_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide_background(slide, C_WHITE)

    # Dot pattern background (decorative circles)
    for i in range(0, 14, 2):
        for j in range(0, 8, 2):
            shape = slide.shapes.add_shape(9, Inches(i * 0.9), Inches(j * 0.9),
                                           Inches(0.15), Inches(0.15))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xE8, 0xDC, 0xF8)
            shape.line.fill.background()

    # Tars logo
    if assets.get("tars_logo") and os.path.exists(assets["tars_logo"]):
        try:
            slide.shapes.add_picture(assets["tars_logo"], Inches(5.2), Inches(0.5),
                                     width=Inches(2.9))
        except Exception:
            add_textbox(slide, "≡TARS", Inches(5.2), Inches(0.5), Inches(3), Inches(1),
                       36, bold=True, color=C_TARS_PURPLE, align=PP_ALIGN.CENTER)
    else:
        add_textbox(slide, "≡TARS", Inches(5.2), Inches(0.5), Inches(3), Inches(1),
                   36, bold=True, color=C_TARS_PURPLE, align=PP_ALIGN.CENTER)

    # Divider line
    line = slide.shapes.add_shape(9, Inches(2.5), Inches(1.7), Inches(8.3), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()

    # Client logo
    if assets.get("client_logo") and os.path.exists(assets["client_logo"]):
        try:
            slide.shapes.add_picture(assets["client_logo"], Inches(4.2), Inches(2.0),
                                     width=Inches(5.0))
        except Exception:
            add_textbox(slide, slide_data.get("client_name", "Client"),
                       Inches(2), Inches(2.2), Inches(9.3), Inches(1.5),
                       48, bold=True, color=C_TARS_PURPLE, align=PP_ALIGN.CENTER)
    else:
        add_textbox(slide, slide_data.get("client_name", "Client"),
                   Inches(2), Inches(2.2), Inches(9.3), Inches(1.5),
                   48, bold=True, color=C_TARS_PURPLE, align=PP_ALIGN.CENTER)

    # Business Review title
    add_textbox(slide, slide_data.get("title", "Business Review"),
               Inches(2), Inches(4.2), Inches(9.3), Inches(0.9),
               36, bold=True, color=C_TARS_PURPLE, align=PP_ALIGN.CENTER)

    # Subtitle (period)
    add_textbox(slide, slide_data.get("subtitle", ""),
               Inches(2), Inches(5.0), Inches(9.3), Inches(0.7),
               20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


def build_overview_stats_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, C_WHITE)
    add_slide_header(prs, slide, slide_data.get("title", "Performance Overview"), assets=assets)

    stats = slide_data.get("stats", [])
    insights = slide_data.get("insights", [])
    period = slide_data.get("period", "")

    if period:
        add_textbox(slide, f"Period: {period}", Inches(0.3), Inches(1.6), Inches(5), Inches(0.4),
                   12, color=C_MID_TEXT)

    # Stat cards - up to 4 in a row
    card_w = Inches(3.0)
    card_h = Inches(1.6)
    start_x = Inches(0.25)
    card_y = Inches(2.0)
    gap = Inches(0.2)

    for i, stat in enumerate(stats[:4]):
        cx = start_x + i * (card_w + gap)
        # Card background
        card = slide.shapes.add_shape(1, cx, card_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = C_LIGHT_PURPLE
        card.line.color.rgb = RGBColor(0xD0, 0xC0, 0xF0)
        card.line.width = Pt(1)

        # Top accent bar
        accent = slide.shapes.add_shape(1, cx, card_y, card_w, Inches(0.07))
        accent.fill.solid()
        accent.fill.fore_color.rgb = C_ACCENT
        accent.line.fill.background()

        # Big number
        add_textbox(slide, stat.get("value", ""),
                   cx + Inches(0.15), card_y + Inches(0.15), card_w - Inches(0.3), Inches(0.8),
                   36, bold=True, color=C_TARS_PURPLE, align=PP_ALIGN.CENTER)

        # Label
        add_textbox(slide, stat.get("label", ""),
                   cx + Inches(0.1), card_y + Inches(0.9), card_w - Inches(0.2), Inches(0.4),
                   13, bold=False, color=C_MID_TEXT, align=PP_ALIGN.CENTER)

        # Note (optional)
        if stat.get("note"):
            add_textbox(slide, stat["note"],
                       cx + Inches(0.1), card_y + Inches(1.25), card_w - Inches(0.2), Inches(0.3),
                       10, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Second row of stats if more than 4
    if len(stats) > 4:
        card_y2 = card_y + card_h + Inches(0.25)
        for i, stat in enumerate(stats[4:8]):
            cx = start_x + i * (card_w + gap)
            card = slide.shapes.add_shape(1, cx, card_y2, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = C_LIGHT_PURPLE
            card.line.color.rgb = RGBColor(0xD0, 0xC0, 0xF0)
            card.line.width = Pt(1)
            accent = slide.shapes.add_shape(1, cx, card_y2, card_w, Inches(0.07))
            accent.fill.solid()
            accent.fill.fore_color.rgb = C_ACCENT
            accent.line.fill.background()
            add_textbox(slide, stat.get("value", ""),
                       cx + Inches(0.15), card_y2 + Inches(0.15), card_w - Inches(0.3), Inches(0.8),
                       36, bold=True, color=C_TARS_PURPLE, align=PP_ALIGN.CENTER)
            add_textbox(slide, stat.get("label", ""),
                       cx + Inches(0.1), card_y2 + Inches(0.9), card_w - Inches(0.2), Inches(0.4),
                       13, color=C_MID_TEXT, align=PP_ALIGN.CENTER)

    # Insights panel on right if space, else below
    if insights:
        insight_y = Inches(2.0)
        insight_x = Inches(12.8) - Inches(0.3)
        # If 4 cards used full width, put insights below
        insights_x = Inches(0.3)
        insights_y = card_y + card_h + Inches(0.5) if len(stats) <= 4 else Inches(6.2)

        for i, insight in enumerate(insights[:4]):
            bullet_x = insights_x
            bullet_y = insights_y + i * Inches(0.45)
            # Bullet dot
            dot = slide.shapes.add_shape(9, bullet_x, bullet_y + Inches(0.12),
                                         Inches(0.12), Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = C_ACCENT
            dot.line.fill.background()
            add_textbox(slide, insight,
                       bullet_x + Inches(0.2), bullet_y, Inches(12.8), Inches(0.4),
                       12, color=C_MID_TEXT)


def build_monthly_trend_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, C_WHITE)
    add_slide_header(prs, slide, slide_data.get("title", "Visits & Conversations - Over the Months"), assets=assets)

    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])
    insights = slide_data.get("insights", [])

    if not headers or not rows:
        add_textbox(slide, "No trend data available.", Inches(0.5), Inches(2), Inches(10), Inches(1), 14)
        return

    # Table
    table_top = Inches(1.65)
    table_left = Inches(0.25)
    table_w = Inches(12.8)

    col_count = len(headers)
    row_count = len(rows) + 1  # +1 for header

    table = slide.shapes.add_table(row_count, col_count, table_left, table_top,
                                   table_w, Inches(0.42) * row_count).table

    col_w = table_w / col_count
    for i in range(col_count):
        table.columns[i].width = int(col_w)

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TABLE_HEADER
        tf = cell.text_frame
        tf.text = h
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = C_WHITE
        run.font.name = "Calibri"

    # Data rows
    for i, row in enumerate(rows):
        bg = C_TABLE_ALT if i % 2 == 0 else C_WHITE
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.text = str(val)
            run.font.size = Pt(11)
            run.font.color.rgb = C_DARK_TEXT
            run.font.name = "Calibri"
            if j == 0:
                run.font.bold = True

    # Insights below table
    insight_y = table_top + Inches(0.45) * row_count + Inches(0.2)
    for i, insight in enumerate(insights[:3]):
        dot = slide.shapes.add_shape(9, Inches(0.3), insight_y + i * Inches(0.45) + Inches(0.12),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_ACCENT
        dot.line.fill.background()
        add_textbox(slide, insight, Inches(0.5), insight_y + i * Inches(0.45),
                   Inches(12.5), Inches(0.42), 12, color=C_MID_TEXT)


def build_generic_table_slide(prs, slide_data, assets):
    """Generic slide with a table and bullet insights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, C_WHITE)
    add_slide_header(prs, slide, slide_data.get("title", "Analysis"), assets=assets)

    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])
    insights = slide_data.get("insights", [])
    description = slide_data.get("description", "")

    if description:
        add_textbox(slide, description, Inches(0.3), Inches(1.6), Inches(12.7), Inches(0.45),
                   12, italic=True, color=C_MID_TEXT)

    table_y = Inches(2.1) if description else Inches(1.65)

    if headers and rows:
        col_count = len(headers)
        row_count = len(rows) + 1
        table_w = Inches(12.8)
        table = slide.shapes.add_table(row_count, col_count, Inches(0.25), table_y,
                                       table_w, Inches(0.42) * row_count).table
        col_w = table_w / col_count
        for i in range(col_count):
            table.columns[i].width = int(col_w)

        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_TABLE_HEADER
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.text = h
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = C_WHITE
            run.font.name = "Calibri"

        for i, row in enumerate(rows):
            bg = C_TABLE_ALT if i % 2 == 0 else C_WHITE
            is_total = str(row[0]).lower() == "total" if row else False
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE0, 0xD0, 0xF5) if is_total else bg
                tf = cell.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                run = p.runs[0] if p.runs else p.add_run()
                run.text = str(val)
                run.font.size = Pt(11)
                run.font.bold = is_total
                run.font.color.rgb = C_DARK_TEXT
                run.font.name = "Calibri"

        insight_y = table_y + Inches(0.42) * row_count + Inches(0.2)
    else:
        insight_y = table_y

    for i, insight in enumerate(insights[:4]):
        if insight_y + i * Inches(0.45) > Inches(7.2):
            break
        dot = slide.shapes.add_shape(9, Inches(0.3),
                                     insight_y + i * Inches(0.45) + Inches(0.12),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_ACCENT
        dot.line.fill.background()
        add_textbox(slide, insight, Inches(0.5), insight_y + i * Inches(0.45),
                   Inches(12.5), Inches(0.42), 12, color=C_MID_TEXT)


def build_insights_summary_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, C_WHITE)
    add_slide_header(prs, slide, slide_data.get("title", "Key Insights"), assets=assets)

    insights = slide_data.get("insights", [])

    for i, insight in enumerate(insights[:8]):
        y = Inches(1.75) + i * Inches(0.65)
        if y > Inches(7.0):
            break
        # Numbered circle
        circle = slide.shapes.add_shape(9, Inches(0.3), y, Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = C_ACCENT
        circle.line.fill.background()
        add_textbox(slide, str(i + 1), Inches(0.3), y - Inches(0.02), Inches(0.4), Inches(0.44),
                   13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, insight, Inches(0.85), y, Inches(12.2), Inches(0.55),
                   13, color=C_DARK_TEXT)


def build_next_steps_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, C_WHITE)
    add_slide_header(prs, slide, slide_data.get("title", "Next Steps for Improvements"), assets=assets)

    items = slide_data.get("items", [])

    for i, item in enumerate(items[:6]):
        y = Inches(1.75) + i * Inches(0.8)
        if y > Inches(7.0):
            break
        # Arrow shape
        arrow = slide.shapes.add_shape(1, Inches(0.3), y + Inches(0.1),
                                       Inches(0.35), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = C_TARS_PURPLE
        arrow.line.fill.background()
        add_textbox(slide, "→", Inches(0.3), y + Inches(0.03), Inches(0.35), Inches(0.4),
                   14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, item, Inches(0.8), y, Inches(12.2), Inches(0.65),
                   14, color=C_DARK_TEXT)


def build_contact_slide(prs, slide_data, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, C_WHITE)

    # Full purple top half
    add_rect(slide, Inches(0), Inches(0), W, Inches(3.0), C_TARS_PURPLE)

    # GET IN TOUCH pink box
    add_rect(slide, Inches(3.5), Inches(0.9), Inches(6.3), Inches(0.7), C_PINK_BG)
    add_rect(slide, Inches(3.5), Inches(0.9), Inches(0.08), Inches(0.7), C_ACCENT)
    add_textbox(slide, "GET IN TOUCH WITH US", Inches(3.7), Inches(0.92), Inches(6), Inches(0.6),
               20, bold=True, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

    # Tars logo on purple
    if assets.get("tars_logo") and os.path.exists(assets["tars_logo"]):
        try:
            slide.shapes.add_picture(assets["tars_logo"], Inches(5.2), Inches(0.1), width=Inches(3.0))
        except Exception:
            add_textbox(slide, "≡TARS", Inches(5.5), Inches(0.15), Inches(2.5), Inches(0.7),
                       24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Contact details
    contacts = [
        ("EMAIL", slide_data.get("email", "team@hellotars.com")),
        ("SOCIAL", slide_data.get("social", "@hellotars.ai")),
        ("WEB", slide_data.get("website", "www.hellotars.com")),
    ]
    for i, (label, val) in enumerate(contacts):
        y = Inches(3.3) + i * Inches(0.9)
        # Label pill
        pill = slide.shapes.add_shape(1, Inches(2.3), y + Inches(0.05), Inches(1.1), Inches(0.45))
        pill.fill.solid()
        pill.fill.fore_color.rgb = C_ACCENT
        pill.line.fill.background()
        add_textbox(slide, label, Inches(2.3), y + Inches(0.05), Inches(1.1), Inches(0.45),
                   10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, val, Inches(3.6), y + Inches(0.05), Inches(6), Inches(0.5),
                   16, bold=True, color=C_TARS_PURPLE)

    # Client logo
    if assets.get("client_logo") and os.path.exists(assets["client_logo"]):
        try:
            slide.shapes.add_picture(assets["client_logo"], Inches(8.5), Inches(3.3),
                                     width=Inches(4.0))
        except Exception:
            pass


# ─── MAIN GENERATOR ───────────────────────────────────────────────────────────

SLIDE_BUILDERS = {
    "cover": build_cover_slide,
    "overview_stats": build_overview_stats_slide,
    "monthly_trend": build_monthly_trend_slide,
    "user_journey": build_generic_table_slide,
    "conversion_funnel": build_generic_table_slide,
    "device_breakdown": build_generic_table_slide,
    "gambit_analysis": build_generic_table_slide,
    "insights_summary": build_insights_summary_slide,
    "next_steps": build_next_steps_slide,
    "contact": build_contact_slide,
}


def generate_pptx(slide_plan: dict, assets: dict, output_path: str):
    """Generate the full PPTX from Claude's slide plan."""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Remove default slide layouts we don't need - keep blank (index 6)
    slides_data = slide_plan.get("slides", [])

    print(f"  🎨 Generating {len(slides_data)} slides...")

    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type", "generic")
        builder = SLIDE_BUILDERS.get(slide_type, build_generic_table_slide)

        # Inject client name into cover slide
        if slide_type == "cover":
            slide_data["client_name"] = assets.get("client_name", "")

        try:
            builder(prs, slide_data, assets)
            print(f"    ✓ Slide {i+1}: {slide_type} - '{slide_data.get('title', slide_type)}'")
        except Exception as e:
            print(f"    ⚠️  Slide {i+1} ({slide_type}) had an issue: {e}")
            # Add a fallback blank slide so we don't crash
            fallback = prs.slides.add_slide(prs.slide_layouts[6])
            slide_background(fallback, C_WHITE)
            add_textbox(fallback, f"[{slide_type}] - {slide_data.get('title', '')}",
                       Inches(1), Inches(2), Inches(11), Inches(2), 24, bold=True)

    prs.save(output_path)
    print(f"  💾 Saved PPTX: {output_path}")
    return output_path
